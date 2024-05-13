import json
from pptx import Presentation
import requests
import base64
from io import BytesIO
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt , Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE


def parse_css_color(css_color):
    # Handles hex, RGB, and RGBA (returns RGBColor or None for transparent)

    if css_color.startswith('#'):
        r, g, b = int(css_color[1:3], 16), int(css_color[3:5], 16), int(css_color[5:7], 16)
    elif css_color.startswith('rgb'):
        parts = css_color[css_color.find('(')+1:css_color.find(')')].split(',')
        r, g, b = int(parts[0].strip()), int(parts[1].strip()), int(parts[2].strip())
        if 'rgba' in css_color and len(parts) == 4 and float(parts[3].strip()) == 0:
            return RGBColor(255, 255, 255) # Treat fully transparent as no color
    else:
        return RGBColor(255, 255, 255)  # Default to white for unrecognized formats
    return RGBColor(r, g, b)


def process_div(slide, content):
    try:
        x_inches = Inches(content['position']['x'] / 96)
        y_inches = Inches(content['position']['y'] / 96)
        width_inches = Inches(content['size']['width'] / 96)
        height_inches = Inches(content['size']['height'] / 96)



        if content['type'] == 'text':
            tx_box = slide.shapes.add_textbox(x_inches, y_inches, width_inches, height_inches)
            tf = tx_box.text_frame
            p = tf.paragraphs[0]
            p.text = content.get('text', '')
            run = p.add_run()
            run.font.size = Pt(int(content.get('styles', {}).get('fontSize', '18').replace("px", "")))  # Default 18pt
            
            color = parse_css_color(content.get('styles', {}).get('color', 'black'))
            if color:
                run.font.color.rgb = color

            # Alignment, bold, and other styles as needed
        elif content['type'] == 'container':
            # If the div serves as a container with background color, consider creating a shape
            # with no text but only the fill color, unless the color is transparent
            fill_color = parse_css_color(content.get('styles', {}).get('backgroundColor', 'transparent'))
            if fill_color:
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, x_inches, y_inches, width_inches, height_inches
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_color
            # Additional handling for div as a generic container
    except Exception as e:
        print("Error", e)

    return None





def process_img(slide, content):
    """
    Adds an image to a slide.
    
    :param slide: The slide to add the image to.
    :param content: A dictionary containing the image's source, position, and size.
    """
    # Convert pixels to inches (assuming 96 DPI)
    x_inches = Inches(content['position']['x'] / 96)
    y_inches = Inches(content['position']['y'] / 96)
    width_inches = Inches(content['size']['width'] / 96)
    height_inches = Inches(content['size']['height'] / 96)

    # Check if the image source is a URL or a local path
    if content['src'].startswith('http'):
        # Image is a URL, download it
        try:
            response = requests.get(content['src'])
            response.raise_for_status()  # Raises HTTPError for bad responses
            image_stream = BytesIO(response.content)
            slide.shapes.add_picture(image_stream, x_inches, y_inches, width_inches, height_inches)
        except requests.HTTPError as e:
            print(f"Failed to download image {content['src']}: {e}")
            return None
    else:
        # Image is a local file
        try:
            print(content['src'].split(","))
            slide.shapes.add_picture(BytesIO(base64.b64decode(content['src'].split(",")[1])), x_inches, y_inches, width_inches, height_inches)
        except FileNotFoundError:
            print(f"File not found: {content['src']}")
            return None
    return None


def download_image(url):
    """Download an image from a URL and return a BytesIO object."""
    response = requests.get(url)
    # Ensure the request was successful
    # print(response.headers)
    
    response.raise_for_status()
        
    return BytesIO(response.content)


def css_font_size_to_pt(css_size):
    # Assuming the font size is provided in pixels, convert to points (1pt = 1.333px)
    if css_size.endswith('px'):
        return Pt(int(css_size[:-2]) * 0.75)  # Simple conversion, adjust as needed
    return Pt(18)  # Default size

def process_text(slide, content):
    x_inches = Inches(content['position']['x'] / 96)  # Convert pixels to inches
    y_inches = Inches(content['position']['y'] / 96)
    width_inches = Inches(content['size']['width'] / 96)
    height_inches = Inches(content['size']['height'] / 96)

    tx_box = slide.shapes.add_textbox(x_inches, y_inches, width_inches, height_inches)
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    p = tf.paragraphs[0]
    p.text = content['text']

    # Font size
    p.font.size = css_font_size_to_pt(content['styles'].get('fontSize', 16))

    # Font color
    p.font.color.rgb = parse_css_color(content['styles'].get('color', '#000000'))

    # Text alignment
    textAlign = content['styles'].get('textAlign', 'left')
    if textAlign == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif textAlign == 'right':
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    # Check if font style includes bold or italic
    font_style = content['styles'].get('fontStyle', '')
    if 'bold' in font_style:
        p.font.bold = True
    if 'italic' in font_style:
        p.font.italic = True

    # Adjust text to fit
    # adjust_text_fit(tf, width_inches, height_inches)
    return None

def adjust_text_fit(text_frame, width, height):
    """
    Adjust the text size to ensure it fits within the specified dimensions.
    The adjustment is heuristic: reducing font size if the text is too long.
    """
    estimated_chars_per_line = width.inches * 2  # Rough estimate: number of characters per line
    max_lines = height.inches / 0.2  # Estimate max number of lines
    max_chars = estimated_chars_per_line * max_lines

    # Reduce font size if text might overflow
    while len(text_frame.text) > max_chars:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                try:
                    if run.font.size.pt > 10:  # Minimum font size threshold
                        run.font.size = Pt(run.font.size.pt - 1)
                    else:
                        return  # Stop if minimum size reached
                except:
                    pass
        try:
            max_chars = estimated_chars_per_line * (height.inches / (0.2 + (0.1 * text_frame.paragraphs.count)))
        except:
            max_chars=100
    return None




def apply_axis_styles(axis, axis_options):
    """Apply styles to chart axes (category or value)."""
    # Axis line color
    axis.format.line.color.rgb = parse_css_color(axis_options['lineColor'])
    # Tick marks color
    # axis.major_tick_mark.format.line.color.rgb = parse_css_color(axis_options['tickColor'])
    # Labels font style
    print(Pt(int(float(axis_options['labels']['fontSize'].replace('em', '')) * 12)))
    axis.tick_labels.font.color.rgb = parse_css_color(axis_options['labels']['color'])
    axis.tick_labels.font.size = Pt(int(float(axis_options['labels']['fontSize'].replace('em', '')) * 12))  # Assuming 'em' units; adjust if necessary

def process_chart(slide, content):
    try:
        # Assuming 96 DPI for now, but consider this might need adjustment based on actual data
        dpi = 96

        # Convert pixel dimensions to inches
        x_inches = Inches(content['position']['x'] / dpi)
        y_inches = Inches(content['position']['y'] / dpi)
        width_inches = Inches(content['size']['width'] / dpi)
        height_inches = Inches(content['size']['height'] / dpi)

        # Debugging output
        print(f"Chart dimensions (inches): x={x_inches}, y={y_inches}, width={width_inches}, height={height_inches}")

        chart_data = CategoryChartData()
        chart_data.categories = content['xAxisCategories']

        for series in content['series']:
            chart_data.add_series(series['name'], tuple([point['y'] for point in series['data']])) 

        chart_type = XL_CHART_TYPE.LINE

        chart = slide.shapes.add_chart(
            chart_type, x_inches, y_inches, width_inches, height_inches, chart_data
        ).chart

        chart.has_title = True
        chart.chart_title.text = content['title']

        # Apply series colors
        for idx, series in enumerate(content['series']):
            chart.series[idx].format.fill.solid()
            chart.series[idx].format.line.color.rgb = parse_css_color(series['color'])

        # Additional customizations as necessary
    except Exception as e:
        print(f"Exception in process_chart: {e}")

    return None


# Style function to apply borders
def set_cell_border(cell,borderTop,borderBottom,borderLeft,borderRight,border_color="#000000" ):
    colors = parse_css_color(border_color)  # Convert hex color to RGBColor object
    
    # Setting borders
    cell.borders.top.color = colors
    cell.borders.top.width = css_font_size_to_pt(borderTop)
    
    cell.borders.bottom.color = colors
    cell.borders.bottom.width = css_font_size_to_pt(borderBottom)
    
    cell.borders.left.color = colors
    cell.borders.left.width = css_font_size_to_pt(borderLeft)
    
    cell.borders.right.color = colors
    cell.borders.right.width = css_font_size_to_pt(borderRight)
    return none


def process_table(slide, content):
    x_inches = Inches(content['position']['x'] / 96)
    y_inches = Inches(content['position']['y'] / 96)
    total_width_inches = Inches(content['size']['width'] / 96)

    rows = len(content['rows'])
    cols = max(len(row['cells']) for row in content['rows'])

    table = slide.shapes.add_table(rows, cols, x_inches, y_inches, total_width_inches, Inches(0)).table

    # Adjust column widths
    if 'columnWidths' in content:
        column_widths = [Inches(w / 96) for w in content['columnWidths']]
        for i, width in enumerate(column_widths):
            if i < len(table.columns):
                table.columns[i].width = width


    # Set cell content and styles
    for i, row in enumerate(content['rows']):
        for j, cell in enumerate(row['cells']):
            cell_obj = table.cell(i, j)
            cell_obj.text = cell['text']
            # set_cell_border(cell_obj,cell['style']['borderTop'],cell['style']['borderBottom'],cell['style']['borderLeft'],cell['style']['borderRight'],cell['style']['color'])
            paragraph = cell_obj.text_frame.paragraphs[0]
            if 'fontSize' in cell['style']:
                paragraph.font.size = css_font_size_to_pt(cell['style']['fontSize'])
            if 'color' in cell['style']:
                paragraph.font.color.rgb = parse_css_color(cell['style']['color'])
            if 'fontStyle' in cell['style'] and cell['style']['fontStyle'].lower()=='italic':
                paragraph.font.italic=True
            if 'fontWeight' in cell['style']:
                try:
                    fontWeight = int(cell['style']['fontWeight'])
                    paragraph.font.bold = fontWeight >= 600
                except ValueError:
                    # Handle 'bold' or other non-numeric values
                    paragraph.font.bold = cell['style']['fontWeight'].lower() == 'bold'


            paragraph.alignment = PP_ALIGN.LEFT

            if 'textAlign' in cell['style']:
                text_align = cell['style']['textAlign'].lower()
                if text_align == 'center':
                    paragraph.alignment = PP_ALIGN.CENTER
                elif text_align == 'right':
                    paragraph.alignment = PP_ALIGN.RIGHT
            
            # Apply additional cell styling here as needed
            # Example: Set cell background color
            if 'backgroundColor' in cell['style']:
                cell_obj.fill.solid()
                cell_obj.fill.fore_color.rgb = parse_css_color(cell['style']['backgroundColor'])
                
            # Example: Set border color and size
            # This is more complex and might require manipulating the XML directly
            # Python-pptx doesn't directly support setting cell border styles as of the last update

    return None



def add_slide_with_background(prs, image_stream):
    # Add a blank slide layout (usually at index 6)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    # Add the background image as a full-slide shape
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    try:
        # Reset stream position to the start
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, 0, 0, width=slide_width, height=slide_height)
    except Exception as e:
        print("image exception",e)
        pass
    return slide








def create_presentation_from_json(json_data,pptx):
    presentation = Presentation(pptx)
    
    data = json_data
    print("data")
    slides=presentation.slides
    for index,slide_data in enumerate(data['presentation']['slides']):

        
        # slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Using a blank layout
        presentation.slide_width = Inches(slide_data['width']/96)
        presentation.slide_height = Inches(slide_data['height']/96)
        slide=slides[index]
        # try:
        #     image_stream = download_image(slide_data['backgroundImage'])
        #     slide= add_slide_with_background(presentation, image_stream)
        #     image_stream.close()
        # except:
        #     slide = add_slide_with_background(presentation, None)
        
            

        for content in slide_data['contents']:
            # print(content['tagName'])
            if content['tagName'].lower() == 'img':
                process_img(slide, content)
            elif content['tagName'].lower() == 'p':
                process_text(slide, content)
            elif content['tagName'].lower() == 'h1':
                process_text(slide, content)
            elif content['tagName'].lower() == 'span':
                process_text(slide, content)
            elif content['tagName'].lower() == 'table':
                process_table(slide, content)
            elif content['tagName'].lower() == 'div':
                process_div(slide, content)
            elif content['tagName'].lower() == 'chart':
                process_chart(slide,content)
            # Add elif for other types as needed

    # presentation.save('generated_presentation.pptx')
    print("Presentation generated successfully.")
    return presentation