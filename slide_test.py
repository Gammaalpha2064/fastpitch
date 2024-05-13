import os
import html
from io import BytesIO
import base64
from pptx import Presentation
import lxml.etree as etree
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from flask import session
import time
from google.cloud import storage
import os
import boto3

storage_client = storage.Client()

# Define the bucket name
bucket_name = 'userstemplates'
bucket = storage_client.bucket(bucket_name)


def emu_to_px(emu, dpi=96):
    """
    Convert EMU to Pixels for given DPI.
    """
    return int(emu / 914400 * 96)


def scale_and_position_shapes(shapes, base_width, base_height, current_width, current_height):
    """
    Scale and position PowerPoint shapes to fit in a designated area.
    """
    scale_width = (current_width / base_width)
    scale_height = (current_height / base_height)
   

    # print(scale)
    scaled_shapes=[]
    for shape in shapes:
        # print("shape",shape)
        scaled_shape={
        'width': emu_to_px(shape['width']) * scale_width,
        'height': emu_to_px(shape['height']) * scale_height,
        'left': emu_to_px(shape['left']) * scale_width,
        'top': emu_to_px(shape['top']) * scale_height + shape.get('vertical_offset', 0),
        'is_placeholder': shape['is_placeholder'],
        'text': shape['text'],
        'type': shape['type'],
        'fill_color': shape['fill_color'],
        'has_text_frame': shape['has_text_frame']}
        # print("Scaled shape",scaled_shape)
        scaled_shapes.append(scaled_shape)

    return scaled_shapes




def image_to_html(shape, style):
    """Convert image shape to HTML."""
    image_stream = BytesIO(shape.image.blob)
    image_base64 = base64.b64encode(image_stream.getvalue()).decode('utf-8')
    return f"<img src='data:image/jpeg;base64,{image_base64}' style='{style}' alt='Embedded Image'/>"

def shape_to_html(shape, background_image=None):
    """
    Converts shape data to HTML format with styles.
    """
    scolor = shape.get('fill_color', "ffffff")

    style = f"position:absolute; left:{shape['left']}px; top:{shape['top']}px; width:{shape['width']}px; height:{shape['height']}px; border:1px solid #ccc; background-color:#{scolor};"
    if shape['is_placeholder']:
        text = html.escape(shape['text'])
        return f"""<div class="dropzone" style="{style}">{text}</div>"""
    elif shape['type'] == 'picture':
        return f"""<div style="{style}"></div>"""
    elif shape['has_text_frame']:
        text = html.escape(shape['text'])
        return f"""<div style="{style}">{text}</div>"""
    return ""



def remove_shapes_and_placeholders(pptx_filename, output_folder,clean_pptx_path):
    pptx_filename=pptx_filename.replace('\\','/')
    clean_pptx_path=clean_pptx_path.replace('\\','/')
    # if pptx_filename[0]=='/':
    #     pptx_filename=pptx_filename[1:]
    # if clean_pptx_path[0]=='/':
    #     clean_pptx_path=clean_pptx_path[1:]

    prs = Presentation(pptx_filename)
    for slide in prs.slides:
        for shape in list(slide.shapes):
            # Remove shape from slide
            sp = shape._element
            sp.getparent().remove(sp)
    # Save the modified presentation

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    prs.save(clean_pptx_path)
    time.sleep(3)

    blob_name = 'cleantemplate.pptx'
    # Create a blob object
    blob = bucket.blob(blob_name)
    # Upload the file to Google Cloud Storage
    blob.upload_from_filename(clean_pptx_path)



    






def slides_to_html(pptx_file, output_folder,  background_images_folder=None):
    """
    Processes each slide in a PowerPoint file and converts to HTML files.
    """
    prs = Presentation(pptx_file)
    base_width, base_height = emu_to_px(prs.slide_width), emu_to_px(prs.slide_height)
    current_width, current_height= ((session.get('width','1280')/2)-20), session.get('height','551')*0.6

    print(base_width,base_height,current_width,current_height)

    # Ensure the output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)


    for slide_number, slide in enumerate(prs.slides, start=0):
        shapes = []
        for shape in slide.shapes:
            fill_color = 'ffffff'  # Default fill color
            try:
                if shape.fill:
                    if hasattr(shape.fill, 'solid') and shape.fill.solid():
                        if shape.fill.fore_color and shape.fill.fore_color.type == 1:
                            fill_color = shape.fill.fore_color.rgb
            except:
                pass
            shapes.append({
                'width': shape.width,
                'height': shape.height,
                'left': shape.left,
                'top': shape.top,
                'is_placeholder': shape.is_placeholder,
                'text': shape.text if shape.has_text_frame else '',
                'type': 'picture' if shape.shape_type == 13 else 'text' if shape.has_text_frame else 'other',
                'fill_color': fill_color,
                'has_text_frame': shape.has_text_frame
            })
        
        scaled_shapes = scale_and_position_shapes(shapes, base_width, base_height, current_width, current_height)
        
        slide_html_content = f"""<div class="slide" style="background-image: url('/background_image/slide_{slide_number}.png'); background-size: 100% 100%;">"""
        for shape in scaled_shapes:
            slide_html_content += shape_to_html(shape)
        slide_html_content += "</div>"

        output_html_file = os.path.join(output_folder, f"slide_{slide_number}.html")
        with open(output_html_file, "w", encoding="utf-8") as file:
            file.write(slide_html_content)
        blob_name = f'slides/slide_{slide_number}.html'
        # Create a blob object in the specified bucket
        blob = bucket.blob(blob_name)
        # Upload the file to Google Cloud Storage
        blob.upload_from_filename(output_html_file)
        # print(f"Slide {slide_number} saved as {output_html_file}")

# Usage example (you'll need to replace 'path/to/presentation.pptx' and 'output/html/slides' with your actual paths):




def process_ppts(pptx_filename,TEMP_DIR):

    clean_pptx_path = TEMP_DIR+"cleantemplate.pptx"  # Update to your PPTX file path
    images_folder = TEMP_DIR+"images"  # Update to your desired output folder path
    slides_folder= TEMP_DIR+"slides"
    


    remove_shapes_and_placeholders(pptx_filename,TEMP_DIR,clean_pptx_path)
    # time.sleep(2)
    # ppt_to_images(clean_pptx_path, images_folder)
    # Prefix to list and delete objects
    prefix = 'slides/'

    # List all blobs with the specified prefix
    blobs = bucket.list_blobs(prefix=prefix)

    # Delete each blob in the list
    for blob in blobs:
        blob.delete()
    slides_to_html(pptx_filename, slides_folder,images_folder)
    return None





# pptx_filename="template.pptx"
# TEMP_DIR=os.getcwd()+"/output/"

# process_ppts(pptx_filename,TEMP_DIR)

# Example usage:
# ppt_to_images(r"/temp/tmp4xifhsus/template.pptx", "/temp/tmp4xifhsus")
